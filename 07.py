from pptx import Presentation

# Укажите путь к файлу презентации
file_path = 'C:/git/Blitz-croco-words/src/Osennyaya_igra_11.pptx'
# Открываем файл презентации
prs = Presentation(file_path)

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        print(shape.text)
