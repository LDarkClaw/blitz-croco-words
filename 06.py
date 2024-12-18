from pptx import Presentation
# Укажите путь к файлу презентации
file_path = 'C:/git/Blitz-croco-words/src/Osennyaya_igra_3.pptx'
# Открываем файл презентации
presentation = Presentation(file_path)
# Теперь вы можете работать с объектом presentation
# Например, выводим количество слайдов
print(f'Количество слайдов в презентации: {len(presentation.slides)}')
# Вы можете итерироваться по слайдам
for slide in presentation.slides:
    print(f'Слайд {slide.slide_id}')