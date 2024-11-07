import zipfile
from pptx import Presentation

words = []

def openzip(name:str):
    list_of_pres: list = []
    a:str = ""
    with zipfile.ZipFile(name, 'r') as zipf:
        for item in zipf.infolist():
            list_of_pres.append(item.filename)

    return(list_of_pres)

def parsing(name):
    prs = Presentation(name)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text: str = shape.text
                if ' ' in text or ':' in text or 'БЛИЦ-КРОКОДИЛ' in text:
                    continue
                words.append(text)
    return words

if __name__ == '__main__':
    list_pres:list = openzip("croco-blitz-source.zip")
    result = []
    for i in range(len(list_pres)):
        result += parsing(list_pres[i])

    print(result)
